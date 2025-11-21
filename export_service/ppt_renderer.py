"""PowerPoint rendering for export payloads."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches, Pt

from .models import ExportRequest


class PowerPointRenderer:
    """Render a basic PPTX deck from session content."""

    def render(self, request: ExportRequest, output_dir: Path) -> Path:
        presentation = Presentation()
        self._add_title_slide(presentation, request)
        if request.sections:
            self._add_section_slides(presentation, request.sections)
        if request.tables:
            self._add_table_slides(presentation, request.tables)
        output_path = output_dir / "report.pptx"
        presentation.save(output_path)
        return output_path

    def _add_title_slide(self, presentation: Presentation, request: ExportRequest) -> None:
        layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = request.title
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = request.summary

    def _add_section_slides(self, presentation: Presentation, sections: Iterable) -> None:
        bullet_layout = presentation.slide_layouts[1]
        for section in sections:
            slide = presentation.slides.add_slide(bullet_layout)
            slide.shapes.title.text = section.heading
            # Preserve basic formatting while keeping the slide readable
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.text = section.body
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)

    def _add_table_slides(self, presentation: Presentation, tables: Iterable) -> None:
        table_layout = presentation.slide_layouts[5]
        for table in tables:
            slide = presentation.slides.add_slide(table_layout)
            slide.shapes.title.text = table.name
            data = [list(table.columns)] + [list(row.dict().values()) for row in table.rows]
            column_count = max(1, len(table.columns))
            row_count = max(1, len(data))
            left = Inches(0.5)
            top = Inches(1.75)
            width = Inches(9)
            height = Inches(0.5 + 0.3 * row_count)
            table_shape = slide.shapes.add_table(row_count, column_count, left, top, width, height).table
            for col_idx, col_name in enumerate(table.columns):
                table_shape.cell(0, col_idx).text = str(col_name)
            for row_idx, row in enumerate(table.rows, start=1):
                row_dict = row.dict()
                for col_idx, col_name in enumerate(table.columns):
                    table_shape.cell(row_idx, col_idx).text = str(row_dict.get(col_name, ""))
            for cell in table_shape.iter_cells():
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
