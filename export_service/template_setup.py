"""Utilities that ensure sample templates exist without storing binaries in git."""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt


def ensure_sample_templates(template_dir: Path, assets_dir: Path, force: bool = False) -> None:
    """Create simple templates + logo if they are missing."""

    template_dir.mkdir(parents=True, exist_ok=True)
    assets_dir.mkdir(parents=True, exist_ok=True)

    _create_logo(assets_dir / "logo.png", force=force)
    _create_summary_template(template_dir / "summary_template.docx", force=force)
    _create_full_report_template(template_dir / "full_report_template.docx", force=force)
    _create_spreadsheet_template(template_dir / "summary_template.xlsx", force=force)
    _create_text_template(template_dir / "summary_template.txt", force=force)
    _create_presentation_template(template_dir / "summary_template.pptx", force=force)


def _create_logo(path: Path, force: bool = False) -> None:
    if path.exists() and not force:
        return

    image = Image.new("RGB", (640, 200), color="#0d47a1")
    draw = ImageDraw.Draw(image)
    text = "Open WebUI"
    try:
        font = ImageFont.truetype("DejaVuSans-Bold.ttf", 64)
    except OSError:
        font = ImageFont.load_default()
    bbox = font.getbbox(text)
    text_width = draw.textlength(text, font=font)
    text_height = bbox[3] - bbox[1]
    draw.text(((640 - text_width) / 2, (200 - text_height) / 2), text, font=font, fill="white")
    image.save(path)


def _create_summary_template(path: Path, force: bool = False) -> None:
    if path.exists() and not force:
        return

    doc = Document()
    doc.core_properties.title = "Summary Report"

    logo_para = doc.add_paragraph("{{ logo }}")
    logo_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    doc.add_heading("{{ title }}", level=1)
    doc.add_paragraph("{{ summary }}")

    doc.add_paragraph("{% if sections %}")
    doc.add_paragraph("{% for section in sections %}")
    doc.add_heading("{{ section.heading }}", level=2)
    doc.add_paragraph("{{ section.body }}")
    doc.add_paragraph("{% endfor %}")
    doc.add_paragraph("{% endif %}")

    doc.add_heading("Tables", level=2)
    doc.add_paragraph("{% if tables %}")
    doc.add_paragraph("{{ tables_json }}")
    doc.add_paragraph("{% else %}No tables supplied.{% endif %}")

    doc.save(path)


def _create_spreadsheet_template(path: Path, force: bool = False) -> None:
    if path.exists() and not force:
        return

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = "Title"
    sheet["B1"] = "{{ title }}"
    sheet["A2"] = "Summary"
    sheet["B2"] = "{{ summary }}"
    sheet["A4"] = "Table name"
    sheet["B4"] = "Value"
    sheet["A5"] = "{% for table in tables %}{{ table.name }}{% endfor %}"
    sheet["B5"] = "Rows: {{ table.rows|length }}"
    workbook.save(path)


def _create_text_template(path: Path, force: bool = False) -> None:
    if path.exists() and not force:
        return

    path.write_text(
        "\n".join(
            [
                "{{ title }}",
                "================",
                "{{ summary }}",
                "",
                "{% for section in sections %}",
                "## {{ section.heading }}",
                "{{ section.body }}",
                "",
                "{% endfor %}",
            ]
        ),
        encoding="utf-8",
    )


def _create_presentation_template(path: Path, force: bool = False) -> None:
    if path.exists() and not force:
        return

    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = "{{ title }}"
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "{{ summary }}"

    bullet_layout = presentation.slide_layouts[1]
    section_slide = presentation.slides.add_slide(bullet_layout)
    section_slide.shapes.title.text = "Sections"
    body_shape = section_slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "{% for section in sections %}{{ section.heading }}"
    p = tf.add_paragraph()
    p.text = "{{ section.body }}"
    p.level = 1
    tf.add_paragraph().text = "{% endfor %}"

    table_layout = presentation.slide_layouts[5]
    table_slide = presentation.slides.add_slide(table_layout)
    table_slide.shapes.title.text = "Tables"
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(1.5)
    table_shape = table_slide.shapes.add_table(2, 2, left, top, width, height).table
    table_shape.cell(0, 0).text = "Table"
    table_shape.cell(0, 1).text = "Rows"
    table_shape.cell(1, 0).text = "{% for table in tables %}{{ table.name }}{% endfor %}"
    table_shape.cell(1, 1).text = "{{ table.rows|length }}"

    for cell in table_shape.iter_cells():
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptPt(14)

    presentation.save(path)


def _create_full_report_template(path: Path, force: bool = False) -> None:
    if path.exists() and not force:
        return

    doc = Document()
    doc.core_properties.title = "Full Report"
    doc.add_paragraph("{{ logo }}").alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_heading("{{ title }}", level=1)
    summary_para = doc.add_paragraph("{{ summary }}")
    if summary_para.runs:
        summary_para.runs[0].font.size = Pt(12)

    doc.add_paragraph("{% for section in sections %}")
    doc.add_heading("{{ section.heading }}", level=2)
    doc.add_paragraph("{{ section.body }}")
    doc.add_paragraph("{% endfor %}")

    doc.add_heading("Tables", level=2)
    doc.add_paragraph("{% for table in tables %}")
    doc.add_heading("{{ table.name }}", level=3)
    doc.add_paragraph("Columns: {{ table.columns }}")
    doc.add_paragraph("{% for row in table.rows %}Row: {{ row.__root__ }}{% endfor %}")
    doc.add_paragraph("{% endfor %}")

    footer = doc.sections[0].footer
    if not footer.paragraphs:
        footer.add_paragraph()
    footer.paragraphs[0].text = "Session {{ session_id }} — Generated {{ title }}"

    doc.save(path)


def main() -> None:  # pragma: no cover - CLI helper
    import argparse

    parser = argparse.ArgumentParser(description="Generate sample DOCX templates and assets.")
    parser.add_argument("--template-dir", type=Path, default=Path("templates"))
    parser.add_argument("--assets-dir", type=Path, default=Path("templates/assets"))
    parser.add_argument("--force", action="store_true", help="Overwrite existing files")
    args = parser.parse_args()

    ensure_sample_templates(args.template_dir, args.assets_dir, force=args.force)
    print(f"Templates ready under {args.template_dir}")


if __name__ == "__main__":  # pragma: no cover
    main()
